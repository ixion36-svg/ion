"""Daily SOC Standup API — aggregates health checks, alerts, cases, and log monitoring."""

import asyncio
import logging
import os
from datetime import datetime, timedelta, timezone
from typing import Any, Dict, List, Optional

from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import Response
from pydantic import BaseModel, Field
from sqlalchemy import func

from ion.auth.dependencies import require_permission
from ion.models.user import User

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/daily-standup", tags=["daily-standup"])


# ── Log-source health config (v0.10.18) ───────────────────────────────────
#
# The DC/WEF checks were originally hard-coded to query winlogbeat-* for
# host.hostname.keyword matching the literal substrings DCS / WEF. Real
# deployments name their hosts differently (e.g. `corp-dc-01`, `evtfwd-prd`)
# AND store events in different indices (`logs-windows.*-`, `filebeat-*`,
# custom data streams). These four env vars let operators point the checks
# at their actual data without code edits.
#
# All defaults preserve v0.10.17 behaviour exactly.

def _parse_patterns(raw: str) -> List[str]:
    """Split comma- or whitespace-separated patterns, drop blanks."""
    if not raw:
        return []
    parts: List[str] = []
    for chunk in raw.replace(";", ",").split(","):
        for sub in chunk.strip().split():
            s = sub.strip()
            if s:
                parts.append(s)
    return parts


def _standup_log_index() -> str:
    return os.environ.get("ION_STANDUP_LOG_INDEX", "winlogbeat-*").strip() or "winlogbeat-*"


def _standup_host_field() -> str:
    return (
        os.environ.get("ION_STANDUP_HOST_FIELD", "host.hostname.keyword").strip()
        or "host.hostname.keyword"
    )


def _standup_dcs_patterns() -> List[str]:
    return _parse_patterns(os.environ.get("ION_STANDUP_DCS_HOSTS", "*DCS*"))


def _standup_wef_patterns() -> List[str]:
    return _parse_patterns(os.environ.get("ION_STANDUP_WEF_HOSTS", "*WEF*"))


# ── Internal check helpers ────────────────────────────────────────────────


async def _check_cluster_health() -> Dict[str, Any]:
    """ES cluster health + stats."""

    from ion.services.elasticsearch_service import ElasticsearchService

    es = ElasticsearchService()
    if not es.is_configured:
        return {"status": "not_configured"}
    try:
        # Use a dedicated client to avoid event-loop-closed errors from shared pool
        result = await es._request("GET", "/_cluster/health")
        try:
            stats = await es._request("GET", "/_cluster/stats")
        except Exception:
            stats = {"indices": {}}
        return {
            "status": result.get("status"),  # green/yellow/red
            "number_of_nodes": result.get("number_of_nodes"),
            "active_shards": result.get("active_shards"),
            "relocating_shards": result.get("relocating_shards"),
            "unassigned_shards": result.get("unassigned_shards"),
            "indices_count": stats.get("indices", {}).get("count", 0),
            "docs_count": stats.get("indices", {}).get("docs", {}).get("count", 0),
            "store_size": stats.get("indices", {}).get("store", {}).get("size_in_bytes", 0),
        }
    except Exception as e:
        return {"status": "error", "error": str(e)[:100]}


async def _check_critical_alerts() -> Dict[str, Any]:
    """Critical alerts in the last 24 h.

    v0.19.5: was Critical + High. Daily standup is meant to surface
    "what should ops act on right now" — the High band was diluting
    the focus. Operators can still see High via /alerts. Standup
    stays Critical-only.
    """
    from ion.services.elasticsearch_service import ElasticsearchService

    es = ElasticsearchService()
    if not es.is_configured:
        return {"total": 0, "alerts": []}
    try:
        alerts = await es.get_alerts(hours=24, severity="critical", limit=50)
        return {
            "critical_count": len(alerts),
            "total": len(alerts),
            "alerts": [
                {
                    "id": a.id,
                    "title": a.title,
                    "severity": a.severity,
                    "status": a.status,
                    "host": a.host,
                    "timestamp": a.timestamp.isoformat(),
                    "rule_name": a.rule_name,
                }
                for a in sorted(alerts, key=lambda x: x.timestamp, reverse=True)[:20]
            ],
        }
    except Exception as e:
        return {"total": 0, "error": str(e)[:100]}


async def _check_stale_cases() -> Dict[str, Any]:
    """Open / acknowledged cases older than 24 h."""
    from ion.core.config import get_config
    from ion.models.alert_triage import AlertCase, AlertCaseStatus
    from ion.storage.database import get_engine, get_session_factory

    config = get_config()
    engine = get_engine(config.db_path)
    factory = get_session_factory(engine)
    session = factory()
    try:
        cutoff = datetime.now(timezone.utc) - timedelta(hours=24)
        # Strip timezone for comparison with naive DB timestamps
        cutoff_naive = cutoff.replace(tzinfo=None)
        stale = (
            session.query(AlertCase)
            .filter(
                AlertCase.status.in_([AlertCaseStatus.OPEN, AlertCaseStatus.ACKNOWLEDGED]),
                AlertCase.created_at < cutoff_naive,
            )
            .order_by(AlertCase.created_at.asc())
            .limit(20)
            .all()
        )
        return {
            "count": len(stale),
            "cases": [
                {
                    "id": c.id,
                    "case_number": c.case_number,
                    "title": c.title,
                    "severity": c.severity,
                    "status": c.status.value if hasattr(c.status, "value") else str(c.status),
                    "created_at": c.created_at.isoformat() if c.created_at else None,
                    "assigned_to": (
                        (c.assigned_to.display_name or c.assigned_to.username)
                        if c.assigned_to
                        else "Unassigned"
                    ),
                    "hours_open": (
                        round((datetime.utcnow() - c.created_at).total_seconds() / 3600)
                        if c.created_at
                        else 0
                    ),
                }
                for c in stale
            ],
        }
    finally:
        session.close()


async def _check_open_alerts_30d() -> Dict[str, Any]:
    """Volume of alerts and the share still open across the last 30 days.

    Reads from ION's local AlertTriage table (where every alert ION has
    seen gets a triage row), so this is honest about analyst workload
    even when ES is rotating older alerts out of its hot indices.
    """
    from ion.core.config import get_config
    from ion.models.alert_triage import AlertTriage, AlertTriageStatus
    from ion.storage.database import get_engine, get_session_factory

    config = get_config()
    engine = get_engine(config.db_path)
    factory = get_session_factory(engine)
    session = factory()
    try:
        cutoff = (datetime.now(timezone.utc) - timedelta(days=30)).replace(tzinfo=None)
        total = (
            session.query(AlertTriage)
            .filter(AlertTriage.created_at >= cutoff)
            .count()
        )
        open_count = (
            session.query(AlertTriage)
            .filter(
                AlertTriage.created_at >= cutoff,
                AlertTriage.status == AlertTriageStatus.OPEN,
            )
            .count()
        )
        ack_count = (
            session.query(AlertTriage)
            .filter(
                AlertTriage.created_at >= cutoff,
                AlertTriage.status == AlertTriageStatus.ACKNOWLEDGED,
            )
            .count()
        )
        closed_count = (
            session.query(AlertTriage)
            .filter(
                AlertTriage.created_at >= cutoff,
                AlertTriage.status == AlertTriageStatus.CLOSED,
            )
            .count()
        )
        return {
            "window_days": 30,
            "total": int(total),
            "open": int(open_count),
            "acknowledged": int(ack_count),
            "closed": int(closed_count),
            "still_open_pct": round(open_count * 100 / total) if total else 0,
        }
    except Exception as e:
        return {"window_days": 30, "total": 0, "error": str(e)[:120]}
    finally:
        session.close()


async def _check_case_status_counts() -> Dict[str, Any]:
    """Open / in-progress / closed case counts (all-time + last 7d delta)."""
    from ion.core.config import get_config
    from ion.models.alert_triage import AlertCase
    from ion.storage.database import get_engine, get_session_factory

    config = get_config()
    engine = get_engine(config.db_path)
    factory = get_session_factory(engine)
    session = factory()
    try:
        # All-time totals by status — cheap, cases table is small.
        rows = (
            session.query(AlertCase.status, func.count(AlertCase.id))
            .group_by(AlertCase.status)
            .all()
        )
        by_status: Dict[str, int] = {
            (s.value if hasattr(s, "value") else str(s)): int(n) for s, n in rows
        }
        # Status names ION uses: OPEN, ACKNOWLEDGED (≈ in-progress), CLOSED.
        opened_7d = (
            session.query(AlertCase)
            .filter(AlertCase.created_at >= (datetime.utcnow() - timedelta(days=7)))
            .count()
        )
        closed_7d = (
            session.query(AlertCase)
            .filter(AlertCase.closed_at.isnot(None))
            .filter(AlertCase.closed_at >= (datetime.utcnow() - timedelta(days=7)))
            .count()
        )
        return {
            "open":         int(by_status.get("open", 0)),
            "acknowledged": int(by_status.get("acknowledged", 0)),
            "closed":       int(by_status.get("closed", 0)),
            "total":        int(sum(by_status.values())),
            "opened_last_7d": int(opened_7d),
            "closed_last_7d": int(closed_7d),
        }
    except Exception as e:
        return {"error": str(e)[:120]}
    finally:
        session.close()


async def _check_triage_throughput_24h() -> Dict[str, Any]:
    """Alerts triaged in the last 24h + average time-to-acknowledge.

    "Triaged" = an AlertTriage row that has moved out of OPEN (i.e.
    status is ACKNOWLEDGED or CLOSED) AND was acknowledged within the
    last 24h. We use ``updated_at`` as the proxy for "first analyst
    touch" since AlertTriage doesn't carry a separate ack timestamp.
    Mean time-to-acknowledge = mean(updated_at − created_at) across
    those rows; reported in minutes.
    """
    from ion.core.config import get_config
    from ion.models.alert_triage import AlertTriage, AlertTriageStatus
    from ion.storage.database import get_engine, get_session_factory

    config = get_config()
    engine = get_engine(config.db_path)
    factory = get_session_factory(engine)
    session = factory()
    try:
        cutoff = datetime.utcnow() - timedelta(hours=24)
        triaged = (
            session.query(AlertTriage)
            .filter(
                AlertTriage.status.in_([AlertTriageStatus.ACKNOWLEDGED, AlertTriageStatus.CLOSED]),
                AlertTriage.updated_at >= cutoff,
            )
            .all()
        )
        if not triaged:
            return {
                "triaged_24h": 0,
                "avg_mtta_minutes": None,
                "p50_mtta_minutes": None,
                "p90_mtta_minutes": None,
            }
        deltas: List[float] = []
        for t in triaged:
            if t.created_at and t.updated_at and t.updated_at >= t.created_at:
                deltas.append((t.updated_at - t.created_at).total_seconds() / 60.0)
        if not deltas:
            return {
                "triaged_24h": len(triaged),
                "avg_mtta_minutes": None,
                "p50_mtta_minutes": None,
                "p90_mtta_minutes": None,
            }
        deltas_sorted = sorted(deltas)
        n = len(deltas_sorted)
        avg = sum(deltas_sorted) / n
        p50 = deltas_sorted[n // 2]
        p90 = deltas_sorted[min(n - 1, int(n * 0.9))]
        return {
            "triaged_24h":      len(triaged),
            "avg_mtta_minutes": round(avg, 1),
            "p50_mtta_minutes": round(p50, 1),
            "p90_mtta_minutes": round(p90, 1),
        }
    except Exception as e:
        return {"triaged_24h": 0, "error": str(e)[:120]}
    finally:
        session.close()


async def _check_log_source_health(
    patterns: List[str],
    label: str,
    *,
    index: Optional[str] = None,
    host_field: Optional[str] = None,
) -> Dict[str, Any]:
    """Check `index` for hosts matching any of `patterns`, compare to 7-day average, find gaps.

    v0.10.18: ``patterns`` is now a list (multiple wildcards OR'd together)
    and ``index`` + ``host_field`` are configurable. Returns a ``diag``
    block in the response so the UI can show ``Queried X for Y → N hosts``
    when nothing comes back — the v0.10.17 silent-empty masked field-name
    and index mismatches.
    """
    from ion.services.elasticsearch_service import ElasticsearchService

    index = index or _standup_log_index()
    host_field = host_field or _standup_host_field()
    diag: Dict[str, Any] = {
        "index": index,
        "host_field": host_field,
        "patterns": list(patterns),
        "total_hits": 0,
        "host_count": 0,
    }

    es = ElasticsearchService()
    if not es.is_configured:
        return {"label": label, "status": "not_configured", "diag": diag}

    if not patterns:
        return {
            "label": label,
            "status": "not_configured",
            "error": f"No host patterns configured (set ION_STANDUP_DCS_HOSTS / ION_STANDUP_WEF_HOSTS)",
            "diag": diag,
        }

    # Build a `should`-clause across all supplied patterns. minimum_should_match=1
    # means ANY pattern hit qualifies — operators can safely list "*DCS*", "*DC*"
    # without one swallowing the other.
    should_clauses = [{"wildcard": {host_field: p}} for p in patterns]

    try:
        # Last 24 h event count per host matching pattern
        # v0.15.1: terms.size dropped 500→100. WEF estates rarely have
        # >100 forwarders in scope, and the 500-host fan-out (×24 hourly
        # buckets) was the dominant cost causing timeouts on busy estates.
        # Operators with more hosts than that can override via
        # ION_STANDUP_TERMS_SIZE.
        try:
            terms_size = int(os.environ.get("ION_STANDUP_TERMS_SIZE", "100"))
        except Exception:
            terms_size = 100
        # v0.15.1: heavy aggregations get a 60 s per-request timeout
        # (vs the 30 s ION_ES_TIMEOUT default). Without this override the
        # WEF check on a busy estate raced the global default and timed
        # out before ES could finish the rollup.
        body_24h = {
            "size": 0,
            "query": {
                "bool": {
                    "should": should_clauses,
                    "minimum_should_match": 1,
                    "filter": [
                        {"range": {"@timestamp": {"gte": "now-24h", "lte": "now"}}},
                    ],
                }
            },
            "aggs": {
                "hosts": {
                    "terms": {"field": host_field, "size": terms_size},
                    "aggs": {
                        "hourly": {
                            "date_histogram": {"field": "@timestamp", "fixed_interval": "1h"}
                        }
                    },
                }
            },
        }
        result_24h = await es._request(
            "POST", f"/{index}/_search?ignore_unavailable=true",
            json=body_24h, timeout=60.0,
        )

        # 7-day average for comparison
        body_7d = {
            "size": 0,
            "query": {
                "bool": {
                    "should": should_clauses,
                    "minimum_should_match": 1,
                    "filter": [
                        {"range": {"@timestamp": {"gte": "now-7d", "lte": "now"}}},
                    ],
                }
            },
            "aggs": {
                "hosts": {
                    "terms": {"field": host_field, "size": terms_size},
                }
            },
        }
        result_7d = await es._request(
            "POST", f"/{index}/_search?ignore_unavailable=true",
            json=body_7d, timeout=60.0,
        )

        hosts_24h: Dict[str, Dict[str, Any]] = {}
        gaps: List[Dict[str, Any]] = []
        for bucket in result_24h.get("aggregations", {}).get("hosts", {}).get("buckets", []):
            hostname = bucket["key"]
            count = bucket["doc_count"]
            hourly = bucket.get("hourly", {}).get("buckets", [])
            gap_hours = [h["key_as_string"] for h in hourly if h["doc_count"] == 0]
            hosts_24h[hostname] = {"count_24h": count, "gap_hours": gap_hours}
            if gap_hours:
                gaps.append({"host": hostname, "gap_count": len(gap_hours), "gaps": gap_hours[:5]})

        # 7-day daily averages
        hosts_7d: Dict[str, float] = {}
        for bucket in result_7d.get("aggregations", {}).get("hosts", {}).get("buckets", []):
            hosts_7d[bucket["key"]] = bucket["doc_count"] / 7

        # Per-host summary
        host_summaries: List[Dict[str, Any]] = []
        for hostname, data in sorted(hosts_24h.items()):
            avg_7d = hosts_7d.get(hostname, 0)
            count_24h = data["count_24h"]
            below_avg = count_24h < (avg_7d * 0.7) if avg_7d > 0 else False
            host_summaries.append(
                {
                    "hostname": hostname,
                    "count_24h": count_24h,
                    "avg_7d": round(avg_7d),
                    "below_average": below_avg,
                    "gap_hours": len(data["gap_hours"]),
                    "status": (
                        "critical"
                        if len(data["gap_hours"]) > 4
                        else "warning"
                        if len(data["gap_hours"]) > 0 or below_avg
                        else "ok"
                    ),
                }
            )

        total_24h = result_24h.get("hits", {}).get("total", {}).get("value", 0)
        diag["total_hits"] = total_24h
        diag["host_count"] = len(hosts_24h)
        return {
            "label": label,
            "status": "ok",
            "total_events_24h": total_24h,
            "host_count": len(hosts_24h),
            "hosts_with_gaps": len(gaps),
            "hosts_below_average": sum(1 for h in host_summaries if h["below_average"]),
            "hosts": sorted(host_summaries, key=lambda h: h["status"] != "ok", reverse=True),
            "gaps": gaps,
            "diag": diag,
        }
    except Exception as e:
        # type(e).__name__ in the message so connection-vs-auth-vs-query
        # failures are distinguishable in the UI banner.
        return {
            "label": label,
            "status": "error",
            "error": f"{type(e).__name__}: {str(e)[:180]}",
            "diag": diag,
        }


async def _check_rule_failures() -> Dict[str, Any]:
    """Detection rules that failed execution in the last 24 h.

    Tries multiple index patterns since Kibana versions store execution
    logs differently:
    - .kibana-event-log-* (older Kibana)
    - .internal.alerts-* with execution status fields
    - .kibana-alerting-* (some 8.x versions)
    """
    from ion.services.elasticsearch_service import ElasticsearchService

    es = ElasticsearchService()
    if not es.is_configured:
        return {"count": 0, "rules": []}

    # Try multiple approaches for rule failure detection
    rules: List[Dict[str, Any]] = []

    # Approach 1: Kibana event log
    indices_to_try = [
        ".kibana-event-log-*",
        ".internal.kibana-event-log-*",
    ]
    fields_to_try = [
        "kibana.alert.rule.execution.metrics.execution_status",
        "event.outcome",
    ]

    for index in indices_to_try:
        if rules:
            break
        for status_field in fields_to_try:
            if rules:
                break
            try:
                body = {
                    "size": 0,
                    "query": {
                        "bool": {
                            "must": [
                                {"range": {"@timestamp": {"gte": "now-24h"}}},
                                {"term": {status_field: "failure"}},
                            ]
                        }
                    },
                    "aggs": {
                        "rules": {
                            "terms": {
                                "field": "rule.name",
                                "size": 20,
                                "missing": "unknown",
                            },
                            "aggs": {
                                "last_failure": {"max": {"field": "@timestamp"}},
                            },
                        }
                    },
                }
                result = await es._request(
                    "POST", f"/{index}/_search?ignore_unavailable=true", json=body
                )
                for bucket in result.get("aggregations", {}).get("rules", {}).get("buckets", []):
                    rules.append({
                        "rule_name": bucket["key"],
                        "failure_count": bucket["doc_count"],
                        "last_failure": bucket.get("last_failure", {}).get("value_as_string"),
                    })
            except Exception:
                continue

    return {"count": len(rules), "rules": rules}


# ── Endpoints ─────────────────────────────────────────────────────────────


@router.get("/checks")
async def get_daily_checks(
    current_user: User = Depends(require_permission("alert:read")),
):
    """Aggregate all daily SOC duty checks in a single call."""
    (
        cluster, alerts, cases,
        dc_health, wef_health, rule_failures,
        alerts_30d, case_status, triage_throughput,
    ) = await asyncio.gather(
        _check_cluster_health(),
        _check_critical_alerts(),
        _check_stale_cases(),
        _check_log_source_health(_standup_dcs_patterns(), "Domain Controllers"),
        _check_log_source_health(_standup_wef_patterns(), "Windows Event Forwarding"),
        _check_rule_failures(),
        _check_open_alerts_30d(),
        _check_case_status_counts(),
        _check_triage_throughput_24h(),
        return_exceptions=True,
    )

    def _safe(val: Any) -> Any:
        if isinstance(val, Exception):
            return {"error": str(val)[:100]}
        return val

    return {
        "timestamp": datetime.now(timezone.utc).isoformat(),
        "cluster_health":     _safe(cluster),
        "critical_alerts":    _safe(alerts),
        "stale_cases":        _safe(cases),
        "dc_log_health":      _safe(dc_health),
        "wef_log_health":     _safe(wef_health),
        "rule_failures":      _safe(rule_failures),
        # v0.17.1 additions
        "open_alerts_30d":    _safe(alerts_30d),
        "case_status_counts": _safe(case_status),
        "triage_throughput":  _safe(triage_throughput),
    }


# ── Save standup ──────────────────────────────────────────────────────────


class StandupSaveRequest(BaseModel):
    """Body for both ``/save`` and ``/pdf``.

    v0.15.2: field names match the frontend canonically. Earlier the
    frontend sent ``threat_summary`` / ``servicenow_incidents`` /
    ``signoff_analyst`` / ``signoff_confirmed`` / ``meetings`` while the
    backend expected different names; pydantic silently dropped the
    unmatched fields so the AI summary, ServiceNow notes, sign-off, and
    meetings checklist never made it to the saved doc or the PDF.

    ``model_config = {"extra": "ignore"}`` is the default; we leave it
    so any future frontend-only diagnostic fields don't 422.
    """

    # Canonical fields used by both save & PDF
    servicenow_notes: str = ""
    additional_notes: str = ""
    analyst_name: str = ""
    signed_off: bool = False
    checks_data: dict = Field(default_factory=dict)
    ai_summary: str = ""
    reports_of_interest: list = Field(default_factory=list)
    meetings: dict = Field(default_factory=dict)
    custom_meeting_item: str = ""


# ── HTML rendering (shared by /save + /pdf) ───────────────────────────────


# Stable label map for the meetings checklist (mirrors the data-item
# values in daily_standup.html). New checklist items added in the
# template should be added here too so the saved/printed report names them.
_MEETING_LABELS = {
    "standup":      "Morning standup attended",
    "handover":     "Shift handover reviewed",
    "case_review":  "Case review completed",
    "threat_intel": "Threat intel briefing reviewed",
    "action_items": "Action items from previous day checked",
}


def _esc(s: Any) -> str:
    """HTML-escape, tolerating non-str inputs."""
    import html as _html
    return _html.escape("" if s is None else str(s))


def _render_standup_html(data: "StandupSaveRequest", current_user: "User") -> str:
    """Render the standup as standalone HTML.

    Used by both ``/save`` (stored as the document's ``rendered_content``
    so the document-export-PDF flow produces a sensible report) and
    ``/pdf`` (sent through WeasyPrint inline). v0.15.2.
    """
    today = datetime.now(timezone.utc).strftime("%d %b %Y")
    analyst = data.analyst_name or current_user.display_name or current_user.username
    checks = data.checks_data or {}

    out: List[str] = []
    out.append(
        "<html><head><style>"
        "body { font-family: Helvetica, sans-serif; color: #1a1a2e; padding: 30px;"
        " font-size: 11px; line-height: 1.5; }"
        "h1 { color: #0f172a; font-size: 22px; border-bottom: 2px solid #6de4ff;"
        " padding-bottom: 6px; }"
        "h2 { color: #334155; font-size: 14px; margin-top: 20px; background: #f1f5f9;"
        " padding: 6px 12px; border-radius: 4px; }"
        ".meta { color: #64748b; font-size: 10px; margin-bottom: 16px; }"
        "table { width: 100%; border-collapse: collapse; margin: 8px 0; font-size: 10px; }"
        "th { text-align: left; padding: 4px 8px; background: #f1f5f9;"
        " border: 1px solid #e2e8f0; font-size: 9px; text-transform: uppercase;"
        " color: #64748b; }"
        "td { padding: 4px 8px; border: 1px solid #e2e8f0; }"
        ".status-ok { color: #16a34a; font-weight: bold; }"
        ".status-warning { color: #d97706; font-weight: bold; }"
        ".status-critical { color: #dc2626; font-weight: bold; }"
        ".status-green { color: #16a34a; }"
        ".status-yellow { color: #d97706; }"
        ".status-red { color: #dc2626; }"
        ".section-notes { background: #f8fafc; border: 1px solid #e2e8f0;"
        " border-radius: 4px; padding: 8px 12px; margin: 8px 0; white-space: pre-wrap; }"
        ".meetings-list { list-style: none; padding-left: 0; margin: 6px 0; }"
        ".meetings-list li { padding: 3px 0; }"
        ".meetings-list .done { color: #16a34a; }"
        ".meetings-list .skipped { color: #94a3b8; text-decoration: line-through; }"
        ".reports-list { list-style: none; padding-left: 0; margin: 6px 0; }"
        ".reports-list li { padding: 4px 0; border-bottom: 1px solid #f1f5f9; }"
        ".reports-list a { color: #0ea5e9; text-decoration: none; }"
        ".signoff { margin-top: 30px; padding: 16px; border: 2px solid #e2e8f0;"
        " border-radius: 8px; }"
        ".footer { margin-top: 30px; border-top: 1px solid #e2e8f0; padding-top: 8px;"
        " color: #94a3b8; font-size: 9px; text-align: center; }"
        "</style></head><body>"
    )
    out.append("<h1>Daily SOC Standup Report</h1>")
    signed = "Yes" if data.signed_off else "No"
    out.append(
        f'<div class="meta">Date: {_esc(today)} &nbsp;|&nbsp; '
        f"Duty Analyst: {_esc(analyst)} &nbsp;|&nbsp; Signed Off: {signed}</div>"
    )

    # -- Cluster Health --------------------------------------------------------
    cluster = checks.get("cluster_health") or {}
    if cluster:
        status = cluster.get("status", "unknown")
        status_class = f"status-{status}" if status in ("green", "yellow", "red") else ""
        out.append("<h2>Elasticsearch Cluster Health</h2>")
        out.append(
            "<table><tr><th>Status</th><th>Nodes</th><th>Indices</th>"
            "<th>Shards</th><th>Unassigned</th></tr>"
        )
        out.append(
            f'<tr><td class="{status_class}">{_esc(status).upper()}</td>'
            f'<td>{_esc(cluster.get("number_of_nodes", "?"))}</td>'
            f'<td>{_esc(cluster.get("indices_count", "?"))}</td>'
            f'<td>{_esc(cluster.get("active_shards", "?"))}</td>'
            f'<td>{_esc(cluster.get("unassigned_shards", "?"))}</td></tr></table>'
        )

    # -- Critical Alerts -------------------------------------------------------
    alerts = checks.get("critical_alerts") or {}
    if alerts:
        out.append(
            f'<h2>Critical Alerts (Last 24h) &mdash; '
            f'{alerts.get("critical_count", 0)} Critical</h2>'
        )
        alert_list = alerts.get("alerts", [])
        if alert_list:
            out.append(
                "<table><tr><th>Time</th><th>Severity</th><th>Rule</th>"
                "<th>Host</th><th>Status</th></tr>"
            )
            for a in alert_list[:15]:
                ts = str(a.get("timestamp", ""))[:16].replace("T", " ")
                out.append(
                    f'<tr><td>{_esc(ts)}</td><td>{_esc(a.get("severity", ""))}</td>'
                    f'<td>{_esc(str(a.get("rule_name", ""))[:50])}</td>'
                    f'<td>{_esc(a.get("host", ""))}</td>'
                    f'<td>{_esc(a.get("status", ""))}</td></tr>'
                )
            out.append("</table>")
        else:
            out.append("<p>No critical alerts in the last 24 hours.</p>")

    # -- Open alerts (last 30 days) — v0.17.1 ---------------------------------
    a30 = checks.get("open_alerts_30d") or {}
    if a30 and a30.get("total") is not None:
        out.append(
            f'<h2>Alert Backlog (Last 30 days) &mdash; {a30.get("total", 0)} alerts, '
            f'{a30.get("open", 0)} still open ({a30.get("still_open_pct", 0)}%)</h2>'
        )
        out.append(
            '<table><tr><th>Total</th><th>Open</th><th>Acknowledged</th><th>Closed</th></tr>'
            f'<tr><td>{a30.get("total", 0):,}</td>'
            f'<td class="status-warning">{a30.get("open", 0):,}</td>'
            f'<td>{a30.get("acknowledged", 0):,}</td>'
            f'<td class="status-ok">{a30.get("closed", 0):,}</td></tr></table>'
        )

    # -- Case status counts — v0.17.1 ----------------------------------------
    cs = checks.get("case_status_counts") or {}
    if cs and "open" in cs:
        out.append(f'<h2>Cases &mdash; {cs.get("total", 0)} total</h2>')
        out.append(
            '<table><tr><th>Open</th><th>In Progress (acknowledged)</th><th>Closed</th>'
            '<th>Opened (last 7d)</th><th>Closed (last 7d)</th></tr>'
            f'<tr><td class="status-warning">{cs.get("open", 0):,}</td>'
            f'<td>{cs.get("acknowledged", 0):,}</td>'
            f'<td class="status-ok">{cs.get("closed", 0):,}</td>'
            f'<td>{cs.get("opened_last_7d", 0):,}</td>'
            f'<td>{cs.get("closed_last_7d", 0):,}</td></tr></table>'
        )

    # -- Triage throughput (last 24h) + MTTA — v0.17.1 -----------------------
    tt = checks.get("triage_throughput") or {}
    if tt and tt.get("triaged_24h") is not None:
        avg = tt.get("avg_mtta_minutes")
        p50 = tt.get("p50_mtta_minutes")
        p90 = tt.get("p90_mtta_minutes")
        out.append(
            f'<h2>Triage Throughput (Last 24h) &mdash; {tt.get("triaged_24h", 0)} alerts triaged</h2>'
        )
        out.append(
            '<table><tr><th>Triaged (24h)</th><th>Avg MTTA (min)</th>'
            '<th>p50 MTTA (min)</th><th>p90 MTTA (min)</th></tr>'
            f'<tr><td>{tt.get("triaged_24h", 0):,}</td>'
            f'<td>{("—" if avg is None else f"{avg:.1f}")}</td>'
            f'<td>{("—" if p50 is None else f"{p50:.1f}")}</td>'
            f'<td>{("—" if p90 is None else f"{p90:.1f}")}</td></tr></table>'
        )

    # -- Stale Cases -----------------------------------------------------------
    stale = checks.get("stale_cases") or {}
    if stale:
        out.append(f'<h2>Stale Cases (Open &gt; 24h) &mdash; {stale.get("count", 0)}</h2>')
        case_list = stale.get("cases", [])
        if case_list:
            out.append(
                "<table><tr><th>Case</th><th>Title</th><th>Severity</th>"
                "<th>Assigned</th><th>Hours Open</th></tr>"
            )
            for c in case_list:
                out.append(
                    f'<tr><td>{_esc(c.get("case_number", ""))}</td>'
                    f'<td>{_esc(str(c.get("title", ""))[:40])}</td>'
                    f'<td>{_esc(c.get("severity", ""))}</td>'
                    f'<td>{_esc(c.get("assigned_to", ""))}</td>'
                    f'<td>{_esc(c.get("hours_open", 0))}</td></tr>'
                )
            out.append("</table>")
        else:
            out.append("<p>No stale cases.</p>")

    # -- DC / WEF Log Health ---------------------------------------------------
    for key, title in [
        ("dc_log_health", "Domain Controller Log Health"),
        ("wef_log_health", "WEF Log Health"),
    ]:
        lh = checks.get(key) or {}
        if lh and lh.get("hosts"):
            out.append(
                f"<h2>{title} &mdash; {lh.get('total_events_24h', 0):,} events, "
                f"{lh.get('hosts_with_gaps', 0)} hosts with gaps</h2>"
            )
            out.append(
                "<table><tr><th>Host</th><th>Events (24h)</th><th>7-Day Avg</th>"
                "<th>Below Avg</th><th>Gap Hours</th><th>Status</th></tr>"
            )
            for h in lh.get("hosts", []):
                hstatus = h.get("status", "ok")
                sc = (
                    "status-critical" if hstatus == "critical"
                    else "status-warning" if hstatus == "warning"
                    else "status-ok"
                )
                out.append(
                    f'<tr><td>{_esc(h.get("hostname", ""))}</td>'
                    f'<td>{int(h.get("count_24h", 0)):,}</td>'
                    f'<td>{int(h.get("avg_7d", 0)):,}</td>'
                    f'<td>{"Yes" if h.get("below_average") else "No"}</td>'
                    f'<td>{_esc(h.get("gap_hours", 0))}</td>'
                    f'<td class="{sc}">{_esc(hstatus).upper()}</td></tr>'
                )
            out.append("</table>")

    # -- Rule Failures ---------------------------------------------------------
    rf = checks.get("rule_failures") or {}
    if rf and rf.get("rules"):
        out.append(f'<h2>Rule Failures &mdash; {rf.get("count", 0)} rules</h2>')
        out.append("<table><tr><th>Rule</th><th>Failures</th><th>Last Failure</th></tr>")
        for r in rf.get("rules", []):
            out.append(
                f'<tr><td>{_esc(str(r.get("rule_name", ""))[:50])}</td>'
                f'<td>{_esc(r.get("failure_count", 0))}</td>'
                f'<td>{_esc(str(r.get("last_failure") or "")[:16])}</td></tr>'
            )
        out.append("</table>")

    # -- Threat Landscape Summary (AI) -----------------------------------------
    if data.ai_summary:
        out.append(f'<h2>Threat Landscape Summary</h2>'
                   f'<div class="section-notes">{_esc(data.ai_summary)}</div>')

    # -- Threat Reports of Interest --------------------------------------------
    if data.reports_of_interest:
        out.append('<h2>Threat Reports of Interest</h2>')
        out.append('<ul class="reports-list">')
        for rpt in data.reports_of_interest:
            if not isinstance(rpt, dict):
                continue
            title = rpt.get("title") or rpt.get("name") or "Untitled"
            url = rpt.get("url") or ""
            published = rpt.get("published") or rpt.get("created_at") or ""
            link = f'<a href="{_esc(url)}">{_esc(title)}</a>' if url else _esc(title)
            meta = f' &mdash; <span style="color:#64748b">{_esc(str(published)[:16])}</span>' if published else ""
            out.append(f"<li>{link}{meta}</li>")
        out.append("</ul>")

    # -- ServiceNow Notes ------------------------------------------------------
    if data.servicenow_notes:
        out.append(f'<h2>ServiceNow Incidents</h2>'
                   f'<div class="section-notes">{_esc(data.servicenow_notes)}</div>')

    # -- Daily Meetings checklist ----------------------------------------------
    meetings = data.meetings or {}
    custom_item_text = (data.custom_meeting_item or "").strip()
    has_meeting_data = bool(meetings) or bool(custom_item_text)
    if has_meeting_data:
        out.append('<h2>Daily Meetings</h2><ul class="meetings-list">')
        for key, label in _MEETING_LABELS.items():
            done = bool(meetings.get(key))
            mark = "[x]" if done else "[ ]"
            cls = "done" if done else "skipped"
            out.append(f'<li class="{cls}">{mark} {_esc(label)}</li>')
        if custom_item_text:
            done = bool(meetings.get("custom"))
            mark = "[x]" if done else "[ ]"
            cls = "done" if done else "skipped"
            out.append(f'<li class="{cls}">{mark} {_esc(custom_item_text)}</li>')
        out.append("</ul>")

    # -- Additional notes ------------------------------------------------------
    if data.additional_notes:
        out.append(f'<h2>Additional Notes</h2>'
                   f'<div class="section-notes">{_esc(data.additional_notes)}</div>')

    # -- Sign-off block --------------------------------------------------------
    out.append(
        f'<div class="signoff"><strong>Duty Analyst:</strong> {_esc(analyst)}<br>'
        f'<strong>Signed Off:</strong> {"Yes" if data.signed_off else "No"}<br>'
        f"<strong>Date:</strong> {_esc(today)}</div>"
    )
    out.append(
        f'<div class="footer">Generated by ION &middot; '
        f"Intelligent Operating Network &middot; {_esc(today)}</div>"
    )
    out.append("</body></html>")
    return "".join(out)


@router.post("/save")
async def save_daily_standup(
    data: StandupSaveRequest,
    current_user: User = Depends(require_permission("alert:read")),
):
    """Save the daily standup report as a document."""
    from ion.core.config import get_config
    from ion.models.document import Document
    from ion.storage.database import get_engine, get_session_factory

    config = get_config()
    engine = get_engine(config.db_path)
    factory = get_session_factory(engine)
    session = factory()

    try:
        today = datetime.now(timezone.utc).strftime("%Y-%m-%d")
        doc = Document(
            name=f"Daily Standup \u2014 {today}",
            rendered_content=_render_standup_html(data, current_user),
            status="active",
            output_format="html",
        )
        session.add(doc)
        session.commit()
        return {"ok": True, "document_id": doc.id, "name": doc.name}
    except Exception as e:
        session.rollback()
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        session.close()


# ── PDF export ────────────────────────────────────────────────────────────


@router.post("/pdf")
async def export_standup_pdf(
    data: StandupSaveRequest,
    current_user: User = Depends(require_permission("alert:read")),
):
    """Export the daily standup report as a PDF (falls back to HTML without WeasyPrint).

    v0.15.2: HTML body comes from the shared ``_render_standup_html``
    helper, identical to what ``/save`` writes into the document store.
    """
    html = _render_standup_html(data, current_user)

    try:
        from weasyprint import HTML as WeasyHTML

        pdf_bytes = WeasyHTML(string=html).write_pdf()
        filename = f"Daily-Standup-{datetime.now(timezone.utc).strftime('%Y-%m-%d')}.pdf"
        return Response(
            content=pdf_bytes,
            media_type="application/pdf",
            headers={"Content-Disposition": f'attachment; filename="{filename}"'},
        )
    except (ImportError, OSError):
        from fastapi.responses import HTMLResponse

        return HTMLResponse(content=html)



# ── PPTX export (v0.19.9) ────────────────────────────────────────────────


def _build_standup_pptx(checks: Dict[str, Any]) -> bytes:
    """Build a presentation-mode PowerPoint deck from a /checks payload.

    One slide per panel — mirrors the HTML deck served at
    /daily-standup/slides. Uses dark-on-light defaults so the file
    looks acceptable when projected from a stock PowerPoint installation
    (the live deck uses dark theme; printed/exported decks usually need
    light backgrounds to print legibly).
    """
    from io import BytesIO
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    SLATE_900 = RGBColor(0x0F, 0x17, 0x2A)
    SLATE_500 = RGBColor(0x64, 0x74, 0x8B)
    SLATE_200 = RGBColor(0xE2, 0xE8, 0xF0)
    CORAL = RGBColor(0xF8, 0x71, 0x71)
    AMBER = RGBColor(0xFB, 0xBF, 0x24)
    EMERALD = RGBColor(0x4A, 0xDE, 0x80)
    CYAN = RGBColor(0x38, 0xBD, 0xF8)

    def _box(slide, x, y, w, h):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tb.text_frame.word_wrap = True
        return tb.text_frame

    def _line(tf, text, *, size=18, bold=False, color=SLATE_200, align=None):
        p = tf.add_paragraph()
        p.text = str(text)
        run = p.runs[0] if p.runs else p.add_run()
        if not p.runs:
            run.text = str(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        if align is not None:
            p.alignment = align

    def _eyebrow(slide, label):
        tf = _box(slide, 0.6, 0.4, 12.5, 0.5)
        tf.text = label.upper()
        run = tf.paragraphs[0].runs[0]
        run.font.size = Pt(11)
        run.font.color.rgb = SLATE_500
        run.font.bold = True

    def _title(slide, title, *, color=SLATE_900):
        tf = _box(slide, 0.6, 0.95, 12.5, 1.2)
        tf.text = title
        run = tf.paragraphs[0].runs[0]
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = color

    def _kpi(slide, x, label, value, *, color=SLATE_900, value_size=64):
        tf = _box(slide, x, 2.6, 3.0, 1.6)
        tf.text = label.upper()
        tf.paragraphs[0].runs[0].font.size = Pt(11)
        tf.paragraphs[0].runs[0].font.color.rgb = SLATE_500
        tf.paragraphs[0].runs[0].font.bold = True
        p2 = tf.add_paragraph()
        p2.text = str(value)
        p2.runs[0].font.size = Pt(value_size)
        p2.runs[0].font.bold = True
        p2.runs[0].font.color.rgb = color

    def _table_block(slide, rows, headers, top=4.4, height=2.6):
        if not rows:
            return
        tbl = slide.shapes.add_table(
            rows=len(rows) + 1, cols=len(headers),
            left=Inches(0.6), top=Inches(top),
            width=Inches(12.1), height=Inches(height),
        ).table
        for j, h in enumerate(headers):
            cell = tbl.cell(0, j)
            cell.text = h
            for r in cell.text_frame.paragraphs[0].runs:
                r.font.size = Pt(11)
                r.font.bold = True
                r.font.color.rgb = SLATE_500
        for i, row in enumerate(rows, start=1):
            for j, val in enumerate(row):
                cell = tbl.cell(i, j)
                cell.text = "" if val is None else str(val)
                for r in cell.text_frame.paragraphs[0].runs:
                    r.font.size = Pt(12)
                    r.font.color.rgb = SLATE_900

    # --- Title slide ---
    s = prs.slides.add_slide(blank)
    _eyebrow(s, "Daily SOC Standup")
    ts_raw = checks.get("timestamp", "")
    date_str = ts_raw.split("T")[0] if ts_raw else datetime.now(timezone.utc).strftime("%Y-%m-%d")
    _title(s, date_str)
    tf = _box(s, 0.6, 2.3, 12.5, 0.7)
    _line(tf, f"Generated {ts_raw}", size=16, color=SLATE_500)

    # --- Cluster ---
    c = checks.get("cluster_health") or {}
    s = prs.slides.add_slide(blank)
    _eyebrow(s, "Section 1 · Elasticsearch Health")
    _title(s, f"Cluster {c.get('cluster_name', '—')}")
    status = str(c.get("status", "unknown")).lower()
    cluster_color = EMERALD if status == "green" else AMBER if status == "yellow" else CORAL
    _kpi(s, 0.6, "Status", c.get("status", "—"), color=cluster_color, value_size=56)
    _kpi(s, 4.0, "Indices", c.get("indices_count", "—"))
    _kpi(s, 7.4, "Storage", c.get("store_size", "—"), value_size=44)
    _kpi(s, 10.4, "Unassigned", c.get("unassigned_shards", 0),
         color=CORAL if (c.get("unassigned_shards") or 0) > 0 else EMERALD)

    # --- Critical alerts ---
    a = checks.get("critical_alerts") or {}
    s = prs.slides.add_slide(blank)
    _eyebrow(s, "Section 2 · Critical Alerts (Last 24h)")
    _title(s, "Critical Alerts", color=CORAL)
    _kpi(s, 0.6, "Critical", a.get("critical_count", 0), color=CORAL, value_size=80)
    rows = [
        [
            (r.get("timestamp", "") or "")[:16].replace("T", " "),
            r.get("severity", ""),
            (r.get("rule_name", "") or "")[:60],
            r.get("host", ""),
        ]
        for r in (a.get("alerts") or [])[:8]
    ]
    if rows:
        _table_block(s, rows, ["Time", "Severity", "Rule", "Host"])

    # --- Stale cases ---
    sc = checks.get("stale_cases") or {}
    s = prs.slides.add_slide(blank)
    _eyebrow(s, "Section 3 · Stale Cases")
    _title(s, "Stale Cases")
    cases = sc.get("cases") or []
    _kpi(s, 0.6, "Open > Threshold", sc.get("count", len(cases)),
         color=AMBER if cases else EMERALD, value_size=64)
    rows = [
        [
            r.get("case_number", ""),
            (r.get("title", "") or "")[:60],
            r.get("severity", ""),
            f"{r.get('hours_open', 0)}h",
        ]
        for r in cases[:8]
    ]
    if rows:
        _table_block(s, rows, ["Case", "Title", "Severity", "Open"])

    # --- Backlog ---
    b = checks.get("open_alerts_30d") or {}
    if b.get("total") is not None:
        s = prs.slides.add_slide(blank)
        _eyebrow(s, "Section 4 · Alert Backlog (30 days)")
        _title(s, "Alert Backlog")
        still_open = b.get("still_open_pct", 0)
        backlog_color = CORAL if still_open >= 30 else AMBER if still_open >= 15 else EMERALD
        _kpi(s, 0.6, "Total", f"{b.get('total', 0):,}", value_size=48)
        _kpi(s, 4.0, "Still Open", f"{b.get('open', 0):,}", color=backlog_color, value_size=48)
        _kpi(s, 7.4, "Acknowledged", f"{b.get('acknowledged', 0):,}", color=AMBER, value_size=48)
        _kpi(s, 10.4, "Closed", f"{b.get('closed', 0):,}", color=EMERALD, value_size=48)
        tf = _box(s, 0.6, 5.0, 12.5, 1.0)
        _line(tf, f"{still_open}% of 30-day alerts still open.", size=20, color=SLATE_500)

    # --- Case status ---
    cs = checks.get("case_status_counts") or {}
    s = prs.slides.add_slide(blank)
    _eyebrow(s, "Section 5 · Case Status")
    _title(s, "Cases at a Glance")
    _kpi(s, 0.6, "Open", cs.get("open", 0), color=AMBER, value_size=72)
    _kpi(s, 4.6, "Investigating", cs.get("investigating", 0), value_size=72)
    _kpi(s, 8.8, "Closed (24h)", cs.get("closed_24h", 0), color=EMERALD, value_size=72)

    # --- Log health ---
    dc = checks.get("dc_log_health") or {}
    wef = checks.get("wef_log_health") or {}
    s = prs.slides.add_slide(blank)
    _eyebrow(s, "Section 6 · Log-Source Health")
    _title(s, "Log Sources")
    _kpi(s, 0.6, "Domain Controllers (silent / total)",
         f"{dc.get('silent_count', 0)} / {dc.get('total', 0)}",
         color=CORAL if (dc.get("silent_count") or 0) > 0 else EMERALD,
         value_size=56)
    _kpi(s, 6.6, "WEF (silent / total)",
         f"{wef.get('silent_count', 0)} / {wef.get('total', 0)}",
         color=CORAL if (wef.get("silent_count") or 0) > 0 else EMERALD,
         value_size=56)

    # --- Rule failures ---
    rf = checks.get("rule_failures") or {}
    rules = rf.get("rules") or []
    s = prs.slides.add_slide(blank)
    _eyebrow(s, "Section 7 · Detection Rule Failures")
    _title(s, "Failing Rules")
    _kpi(s, 0.6, "Failing rules", rf.get("count", len(rules)),
         color=CORAL if rules else EMERALD, value_size=72)
    if rules:
        _table_block(s, [[(r.get("name") or "")[:80], r.get("last_status", "")]
                         for r in rules[:8]], ["Rule", "Last Status"])

    # --- AI Threat Summary (v0.19.11) ---
    ai_summary = checks.get("_ai_summary") or ""
    if ai_summary:
        s = prs.slides.add_slide(blank)
        _eyebrow(s, "Section 8 · Threat Landscape")
        _title(s, "AI Threat Summary")
        # Long-form prose — wrap as a single paragraph in a wide text box.
        tf = _box(s, 0.6, 2.4, 12.1, 4.4)
        tf.word_wrap = True
        # Truncate to the first ~1800 chars so it fits one slide; longer
        # output just gets a "(truncated)" suffix. Multi-slide overflow
        # is intentional follow-up work.
        body = ai_summary if len(ai_summary) <= 1800 else ai_summary[:1800].rstrip() + "\n\n…(truncated)"
        for line in body.split("\n"):
            p = tf.paragraphs[0] if not tf.text else tf.add_paragraph()
            if not p.runs:
                p.text = line
            r = p.runs[0]
            r.text = line
            r.font.size = Pt(16)
            r.font.color.rgb = SLATE_900

    # --- Any Other Business (v0.19.11) ---
    aob = checks.get("_aob") or ""
    s = prs.slides.add_slide(blank)
    _eyebrow(s, "Section 9 · Any Other Business")
    _title(s, "AOB")
    if aob:
        tf = _box(s, 0.6, 2.4, 12.1, 4.4)
        tf.word_wrap = True
        for line in aob.split("\n"):
            p = tf.paragraphs[0] if not tf.text else tf.add_paragraph()
            if not p.runs:
                p.text = line
            r = p.runs[0]
            r.text = line
            r.font.size = Pt(18)
            r.font.color.rgb = SLATE_900
    else:
        tf = _box(s, 0.6, 3.5, 12.1, 1.0)
        _line(tf, "No AOB items recorded.", size=20, color=SLATE_500)

    # --- Closing slide ---
    s = prs.slides.add_slide(blank)
    _eyebrow(s, "End of Standup")
    _title(s, "Questions? Action items?")

    out = BytesIO()
    prs.save(out)
    return out.getvalue()


class StandupPptxRequest(BaseModel):
    """v0.19.11: optional AI threat summary + AOB notes for the .pptx
    deck. The slide-deck and live-page download buttons supply these
    from localStorage; server-only callers can omit them and the
    corresponding slides degrade to placeholders."""
    ai_summary: str = ""
    aob: str = ""


async def _gather_standup_checks_for_pptx(extras: Optional[dict] = None) -> Dict[str, Any]:
    cluster, alerts, cases, dc_health, wef_health, rule_failures, alerts_30d, case_status, _ = (
        await asyncio.gather(
            _check_cluster_health(),
            _check_critical_alerts(),
            _check_stale_cases(),
            _check_log_source_health(_standup_dcs_patterns(), "Domain Controllers"),
            _check_log_source_health(_standup_wef_patterns(), "Windows Event Forwarding"),
            _check_rule_failures(),
            _check_open_alerts_30d(),
            _check_case_status_counts(),
            _check_triage_throughput_24h(),
            return_exceptions=True,
        )
    )

    def _safe(val: Any) -> Any:
        return {"error": str(val)[:100]} if isinstance(val, Exception) else val

    out = {
        "timestamp": datetime.now(timezone.utc).isoformat(),
        "cluster_health":     _safe(cluster),
        "critical_alerts":    _safe(alerts),
        "stale_cases":        _safe(cases),
        "dc_log_health":      _safe(dc_health),
        "wef_log_health":     _safe(wef_health),
        "rule_failures":      _safe(rule_failures),
        "open_alerts_30d":    _safe(alerts_30d),
        "case_status_counts": _safe(case_status),
    }
    if extras:
        out["_ai_summary"] = extras.get("ai_summary") or ""
        out["_aob"] = extras.get("aob") or ""
    return out


@router.post("/pptx")
async def export_standup_pptx_post(
    body: StandupPptxRequest,
    current_user: User = Depends(require_permission("alert:read")),
):
    """v0.19.11: enriched PPTX — caller supplies the AI threat summary
    and AOB notes (typically from localStorage on /daily-standup or
    /daily-standup/slides). The same data the live page uses, just
    sealed into a downloadable deck.
    """
    checks = await _gather_standup_checks_for_pptx(
        extras={"ai_summary": body.ai_summary, "aob": body.aob}
    )
    try:
        pptx_bytes = _build_standup_pptx(checks)
    except ImportError:
        raise HTTPException(
            status_code=501,
            detail="python-pptx is not installed in this image; rebuild with the v0.19.9+ dependency manifest.",
        )
    filename = f"Daily-Standup-{datetime.now(timezone.utc).strftime('%Y-%m-%d')}.pptx"
    return Response(
        content=pptx_bytes,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@router.get("/pptx")
async def export_standup_pptx(
    current_user: User = Depends(require_permission("alert:read")),
):
    """v0.19.9: Download the daily standup as a PowerPoint deck.

    Server-side gathers the same checks payload the live page consumes
    and builds a one-slide-per-panel deck via python-pptx. No client-
    side rendering — file lands as an attachment ready for email,
    PowerPoint, Keynote, LibreOffice.

    The GET form omits the AI threat summary + AOB sections (no client
    state to draw on). The POST form (added in v0.19.11) accepts both
    in the body and includes them as their own slides.
    """
    checks = await _gather_standup_checks_for_pptx()
    try:
        pptx_bytes = _build_standup_pptx(checks)
    except ImportError:
        raise HTTPException(
            status_code=501,
            detail="python-pptx is not installed in this image; rebuild with the v0.19.9+ dependency manifest.",
        )
    filename = f"Daily-Standup-{datetime.now(timezone.utc).strftime('%Y-%m-%d')}.pptx"
    return Response(
        content=pptx_bytes,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


# ── Arkime high-risk traffic ──────────────────────────────────────────────


@router.get("/arkime/high-risk-traffic")
async def check_arkime_high_risk(
    current_user: User = Depends(require_permission("alert:read")),
):
    """Check Arkime for ingress/egress traffic to high-risk countries."""
    from ion.services.arkime_service import get_arkime_service

    svc = get_arkime_service()
    if not svc.is_configured:
        return {"configured": False}

    countries = {
        "RU": "Russia",
        "CN": "China",
        "IR": "Iran",
        "KP": "North Korea",
        "SY": "Syria",
    }

    # First get the list of nodes
    headers = await svc._headers()
    client = await svc._client()
    nodes = []
    try:
        stats_resp = await client.get(
            f"{svc.url}/api/stats", headers=headers,
        )
        if stats_resp.status_code == 200:
            content_type = stats_resp.headers.get("content-type", "")
            if "json" in content_type:
                for n in stats_resp.json().get("data", []):
                    nodes.append(n.get("nodeName") or n.get("id") or "unknown")
    except Exception:
        pass
    if not nodes:
        nodes = ["default"]

    start_time = str(int((datetime.now(timezone.utc) - timedelta(hours=24)).timestamp()))
    stop_time = str(int(datetime.now(timezone.utc).timestamp()))

    results: List[Dict[str, Any]] = []
    for code, name in countries.items():
        country_result = {
            "country_code": code,
            "country": name,
            "session_count": 0,
            "nodes": [],
        }

        for node_name in nodes:
            try:
                resp = await client.get(
                    f"{svc.url}/api/sessions",
                    headers=headers,
                    params={
                        "expression": f'country == {code} && node == "{node_name}"',
                        "startTime": start_time,
                        "stopTime": stop_time,
                        "length": "5",
                        "fields": "id,node,srcIp,dstIp,ipProtocol,protocol,bytes,packets,firstPacket,lastPacket",
                    },
                )
                if resp.status_code == 200:
                    content_type = resp.headers.get("content-type", "")
                    if "json" not in content_type:
                        continue
                    body = resp.json()
                    count = body.get("recordsFiltered", 0)
                    sessions = []
                    for s in (body.get("data") or [])[:5]:
                        sessions.append({
                            "src_ip": s.get("srcIp"),
                            "dst_ip": s.get("dstIp"),
                            "protocol": s.get("protocol") or s.get("ipProtocol"),
                            "bytes": s.get("bytes", 0),
                            "packets": s.get("packets", 0),
                            "first_packet": s.get("firstPacket"),
                            "last_packet": s.get("lastPacket"),
                        })
                    if count > 0:
                        country_result["nodes"].append({
                            "node": node_name,
                            "session_count": count,
                            "sessions": sessions,
                        })
                        country_result["session_count"] += count
            except Exception as e:
                country_result.setdefault("errors", []).append(f"{node_name}: {str(e)[:60]}")

        results.append(country_result)

    return {"countries": results, "total": sum(r["session_count"] for r in results)}


@router.get("/arkime/node-stats")
async def arkime_node_stats(
    current_user: User = Depends(require_permission("alert:read")),
):
    """Pull per-node capture stats from Arkime — packets, bytes, sessions, drops."""
    from ion.services.arkime_service import get_arkime_service

    svc = get_arkime_service()
    if not svc.is_configured:
        return {"configured": False}

    try:
        headers = await svc._headers()
        client = await svc._client()

        # Node stats
        resp = await client.get(
            f"{svc.url}/api/stats",
            headers=headers,
        )
        if resp.status_code != 200:
            return {"configured": True, "error": f"HTTP {resp.status_code}"}

        content_type = resp.headers.get("content-type", "")
        if "json" not in content_type:
            return {"configured": True, "error": "Non-JSON response from Arkime"}

        data = resp.json()
        nodes = []
        for node in data.get("data", []):
            nodes.append({
                "name": node.get("nodeName", "?"),
                "id": node.get("id", "?"),
                "cpu": node.get("cpu", 0),
                "memory": node.get("memory", 0),
                "packets_24h": node.get("deltaPackets", 0),
                "bytes_24h": node.get("deltaBytes", 0),
                "sessions_24h": node.get("deltaSessions", 0),
                "dropped": node.get("deltaDropped", 0),
                "es_dropped": node.get("deltaESDropped", 0),
                "overload_dropped": node.get("deltaOverloadDropped", 0),
                "packets_per_sec": node.get("deltaPacketsPerSec", 0),
                "bytes_per_sec": node.get("deltaBytesPerSec", 0),
                "sessions_per_sec": node.get("deltaSessionsPerSec", 0),
                "disk_queue": node.get("diskQueue", 0),
                "close_queue": node.get("closeQueue", 0),
                "free_space_g": node.get("freeSpaceG", 0),
                "current_time": node.get("currentTime", 0),
            })

        # Also grab ES health through Arkime
        es_health = {}
        try:
            es_resp = await client.get(
                f"{svc.url}/api/eshealth",
                headers=headers,
            )
            if es_resp.status_code == 200 and "json" in es_resp.headers.get("content-type", ""):
                es_health = es_resp.json()
        except Exception:
            pass

        return {
            "configured": True,
            "node_count": len(nodes),
            "nodes": nodes,
            "es_health": {
                "status": es_health.get("status", "unknown"),
                "nodes": es_health.get("number_of_nodes", 0),
                "shards": es_health.get("active_shards", 0),
                "unassigned": es_health.get("unassigned_shards", 0),
            },
        }
    except Exception as e:
        return {"configured": True, "error": str(e)[:100]}
